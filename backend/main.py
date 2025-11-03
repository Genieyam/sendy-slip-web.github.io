# main.py ── FastAPI 통합본 (미리보기 + PPT 생성)
from fastapi import FastAPI, UploadFile, File
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Pt
import pandas as pd
import uvicorn
import re
import os
import io
from fastapi.responses import StreamingResponse
app = FastAPI()

@app.get("/")
def root():
    return {"status": "ok", "message": "backend alive"}

# ──────────────────────────────
# CORS 설정 (React 3000/3001 허용)
# ──────────────────────────────
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://localhost:3001",
        "https://sendy-slip-web-github-io.vercel.app"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ──────────────────────────────
# 엑셀 필터 함수
# ──────────────────────────────
def filter_excel_data(df: pd.DataFrame) -> pd.DataFrame:
    """F/I/L 열 키워드 자동 필터"""
    df = df.copy()
    drop_idx = []
    for i, row in df.iterrows():
        fVal = str(row.get("이름", "")).strip().lower()          # F열
        iVal = str(row.get("상차지상세", "")).strip().lower()    # I열
        lVal = str(row.get("하차지상세", "")).strip().lower()    # L열
        if (
            "scm" in fVal
            or "동신프라텍" in iVal or "엔피씨" in iVal or "물류센터" in iVal or "골드라인" in iVal
            or "동신프라텍" in lVal or "엔피씨" in lVal or "물류센터" in lVal or "골드라인" in lVal
        ):
            drop_idx.append(i)
    return df.drop(index=drop_idx).reset_index(drop=True)

# ──────────────────────────────
# (N11)/(N12)/기타 수량 추출
# ──────────────────────────────
def extract_numbers_from_text(text: str):
    text = text or ""
    n11_match = re.search(r"\(N11\)\s*([\d,\.]+)", text)
    n12_match = re.search(r"\(N12\)\s*([\d,\.]+)", text)
    other_match = re.search(r"\([^)]*\)\s*([\d,\.]+)", text)

    n11 = int(float(n11_match.group(1).replace(",", ""))) if n11_match else 0
    n12 = int(float(n12_match.group(1).replace(",", ""))) if n12_match else 0
    other = int(float(other_match.group(1).replace(",", ""))) if other_match and not (n11 or n12) else 0
    return n11, n12, other

# ──────────────────────────────
# 행 → A~K 딕셔너리로 변환
# ──────────────────────────────
def build_record_from_row(row: pd.Series):
    A = str(row.get("상차지상세", "")).strip()
    B = str(row.get("하차지상세", "")).strip()
    C = str(row.get("배차차량정보", "")).strip()
    D_raw = str(row.get("짐내용", "")).strip()
    H = str(row.get("실행일", "")).strip()
    K = str(row.get("번호", ""))  # ✅ 엑셀 B열 "번호"

    # A: (점심시간) 제거 + 괄호 전 텍스트 정리
    A = re.sub(r"\(.*점심시간.*\)", "", A)
    if "(" not in A:
        A = re.sub(r"^[^,/]+[,/]\s*", "", A)

    # D,E,F,G 계산
    n11, n12, other = extract_numbers_from_text(D_raw)
    D = str(n11) if n11 else ""
    E = str(n12) if n12 else ""
    F = str(other) if other else ""
    G = str(sum(x for x in (n11, n12, other) if x)) if (n11 or n12 or other) else ""

    # 날짜 형식
    if re.match(r"\d{4}-\d{2}-\d{2}$", H):
        y, m, d = H.split("-")
        H_fmt = f"{y}              {m}          {d}                                                              {y}              {m}          {d}"
    else:
        H_fmt = H

    # I: 날짜+차량번호 뒷자리
    H_digits = re.sub(r"\D", "", H_fmt)[-4:]
    C_digits = re.sub(r"\D", "", C)[-4:]
    I = f"{H_digits}{C_digits}"

    # J: 고정값
    J = "주식회사 센디"

    return {
        "A": A, "B": B, "C": C, "D": D, "E": E,
        "F": F, "G": G, "H": H_fmt, "I": I,
        "J": J, "K": K
    }

# ──────────────────────────────
# PPT 슬라이드 채우기
# ──────────────────────────────
def fill_slide(slide, rec: dict):
    for shp in slide.shapes:
        if not hasattr(shp, "text_frame") or not shp.text_frame:
            continue
        key = (shp.text or "").strip()
        for k in rec.keys():
            if key == k or k in key:
                shp.text_frame.text = rec.get(k, "")
                for p in shp.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(14)
                break

# ──────────────────────────────
# (1) 미리보기 JSON
# ──────────────────────────────
@app.post("/preview_excel/")
async def preview_excel(file: UploadFile = File(...)):
    contents = await file.read()
    df = pd.read_excel(io.BytesIO(contents))
    df = filter_excel_data(df)
    items = [build_record_from_row(row) for _, row in df.iterrows()]
    return JSONResponse({"items": items})

# ──────────────────────────────
# (2) PPT 생성 및 다운로드
# ──────────────────────────────
@app.post("/upload_excel/")
async def upload_excel(file: UploadFile = File(...)):
    contents = await file.read()
    df = pd.read_excel(io.BytesIO(contents))
    df = filter_excel_data(df)
    df.to_excel("filtered_result.xlsx", index=False)

    template_path = "55.pptx"
    if not os.path.exists(template_path):
        return JSONResponse({"error": "❌ 55.pptx가 backend 폴더에 없습니다."}, status_code=400)

    prs = Presentation(template_path)
    total_slides = len(prs.slides)
    rows = min(len(df), total_slides)

    for i in range(rows):
        rec = build_record_from_row(df.iloc[i])
        fill_slide(prs.slides[i], rec)

    out_path = "SendySlip_Result.pptx"
    prs.save(out_path)
    return FileResponse(out_path, filename="SendySlip_Result.pptx")


@app.post("/filter_excel/")
async def filter_excel(file: UploadFile = File(...)):
    """
    업로드된 엑셀을 필터링하고, 필터된 엑셀 파일을 바로 다운로드로 반환
    """
    contents = await file.read()
    df = pd.read_excel(io.BytesIO(contents))
    df_filtered = filter_excel_data(df)

    # 메모리 버퍼에 엑셀 저장
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df_filtered.to_excel(writer, index=False, sheet_name="filtered")
    output.seek(0)

    headers = {
        "Content-Disposition": 'attachment; filename="filtered_result.xlsx"'
    }
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers=headers
    )
# ──────────────────────────────
if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
